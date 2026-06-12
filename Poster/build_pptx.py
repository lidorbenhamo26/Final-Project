"""Build the editable PPTX twin of the Mission Focus A0 poster (for Canva import).
Same theme/content as MissionFocus_Poster_A0.html, rebuilt as native shapes + text.
Run:  python build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
A = lambda name: os.path.join(HERE, "Assets", name)

# ---------- palette ----------
BG      = RGBColor(0x06, 0x0D, 0x1F)
CARD    = RGBColor(0x0A, 0x14, 0x2C)
CARD2   = RGBColor(0x12, 0x24, 0x48)
EDGE    = RGBColor(0x3E, 0x6F, 0xA0)
INK     = RGBColor(0xE8, 0xF0, 0xFF)
SOFT    = RGBColor(0xDF, 0xE9, 0xFB)
DIM     = RGBColor(0xBC, 0xCD, 0xE8)
CYAN    = RGBColor(0x4F, 0xD8, 0xFF)
CYANSOFT= RGBColor(0x9F, 0xDC, 0xFF)
ORANGE  = RGBColor(0xFF, 0x9B, 0x42)
VIOLET  = RGBColor(0xB1, 0x8C, 0xFF)
GREEN   = RGBColor(0x5E, 0xE6, 0xA8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
NAVY    = RGBColor(0x06, 0x12, 0x1F)

ORB = "Orbitron"
OS_ = "Open Sans"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Mm(841)
prs.slide_height = Mm(1189)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def no_line(shape):
    shape.line.fill.background()


def rect(x, y, w, h, fill, line=None, line_w=0.7, radius=None, dash=None, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Mm(x), Mm(y), Mm(w), Mm(h))
    if radius is not None:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        no_line(s)
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
        if dash:
            s.line.dash_style = dash
    s.shadow.inherit = False
    return s


def text(x, y, w, h, runs_lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.18, space_after=4):
    """runs_lines: list of paragraphs; each paragraph = list of run dicts
    {t, font, size, color, bold, italic, spacing}"""
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for r in para:
            run = p.add_run()
            run.text = r["t"]
            f = run.font
            f.name = r.get("font", OS_)
            f.size = Pt(r.get("size", 18))
            f.color.rgb = r.get("color", SOFT)
            f.bold = r.get("bold", False)
            f.italic = r.get("italic", False)
            f.language_id = MSO_LANGUAGE_ID.ENGLISH_US
    return tb


def R(t, **kw):
    d = {"t": t}; d.update(kw); return d


def sec_title(x, y, w, num, title, title_size=27):
    runs = []
    if num:
        runs.append(R(num + "  ", font=ORB, size=17, color=CYAN, bold=True))
    runs.append(R(title.upper(), font=ORB, size=title_size, color=WHITE, bold=True))
    text(x, y, w, 16, [runs])


# ============ BACKGROUND ============
rect(0, 0, 841, 1189, BG)
# stars
import random
random.seed(26117)
for _ in range(220):
    sx, sy = random.uniform(2, 839), random.uniform(2, 1185)
    sr = random.uniform(0.35, 0.9)
    c = random.choice([WHITE, CYANSOFT, RGBColor(0xBE, 0xE6, 0xFF)])
    st = slide.shapes.add_shape(MSO_SHAPE.OVAL, Mm(sx), Mm(sy), Mm(sr), Mm(sr))
    st.fill.solid(); st.fill.fore_color.rgb = c
    no_line(st); st.shadow.inherit = False

# ============ HEADER ============
# logo chip
rect(22, 16, 96, 40, WHITE, radius=0.18)
slide.shapes.add_picture(A("braude_logo.png"), Mm(29), Mm(21), height=Mm(30))
# badges
rect(548, 18, 162, 15.5, CARD, line=EDGE, radius=0.5)
text(548, 21.4, 162, 10, [[R("CAPSTONE PROJECT — PHASE B", font=ORB, size=15, color=CYANSOFT, bold=True)]], align=PP_ALIGN.CENTER)
rect(716, 18, 103, 15.5, CYAN, radius=0.5)
text(716, 21.4, 103, 10, [[R("TEAM 26-1-D-17", font=ORB, size=15, color=NAVY, bold=True)]], align=PP_ALIGN.CENTER)

# title block
text(22, 64, 600, 60, [[
    R("MISSION ", font=ORB, size=128, color=WHITE, bold=True),
    R("FOCUS", font=ORB, size=128, color=CYAN, bold=True),
]], line_spacing=1.0)
text(22, 120, 600, 20, [[R("A SPACE-STATION COGNITIVE ASSESSMENT GAME", font=OS_, size=37, color=CYANSOFT, bold=True)]])
text(22, 146, 600, 18, [[
    R("Lidor Ben Hamo ", font=OS_, size=33, color=INK, bold=True),
    R("& ", font=OS_, size=33, color=ORANGE, bold=True),
    R("Yahli Rapaport", font=OS_, size=33, color=INK, bold=True),
]])
text(22, 165, 700, 14, [[
    R("Advisor: ", font=OS_, size=24, color=SOFT),
    R("Dr. Moshe Sulamy", font=OS_, size=24, color=WHITE, bold=True),
    R("   ·   Department of Software Engineering & Information Systems", font=OS_, size=24, color=SOFT),
]])

# hero ring + astronaut
ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Mm(648), Mm(28), Mm(138), Mm(138))
ring.fill.solid(); ring.fill.fore_color.rgb = RGBColor(0x0B, 0x1E, 0x3E)
ring.line.color.rgb = EDGE; ring.line.width = Pt(2); ring.shadow.inherit = False
slide.shapes.add_picture(A("char_hero.png"), Mm(642), Mm(40), height=Mm(148))

# ============ COLUMN GEOMETRY ============
LX, CX, RX = 22, 282, 569
LW, CW, RW = 250, 277, 250
TOP, BOTTOM = 200, 1042


def stack(x, w, cards):
    """cards: list of (height, draw_fn(x, y, w, h)). Distributes leftover as equal gaps.
    Never exceeds BOTTOM: gap floors at 3mm and a warning is printed if cards over-fill."""
    total = sum(h for h, _ in cards)
    gaps = len(cards) - 1
    leftover = BOTTOM - TOP - total
    if leftover < 3 * gaps:
        print(f"WARNING: column at x={x} over budget by {3 * gaps - leftover:.0f}mm — shrink card heights")
    gap = max(3, leftover / gaps) if gaps else 0
    y = TOP
    for h, fn in cards:
        rect(x, y, w, h, CARD, line=EDGE, radius=0.045)
        fn(x, y, w, h)
        y += h + gap


PAD = 11  # card padding


def bullets(x, y, w, items, size=20, gap=6, marker="◆", mcolor=CYAN, msize=None):
    paras = []
    for runs in items:
        paras.append([R(marker + "  ", font=OS_, size=msize or size, color=mcolor, bold=True)] + runs)
    text(x, y, w, 200, paras, space_after=gap, line_spacing=1.2)


# ---------- LEFT ----------
def bg_need(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "01", "Background & Need")
    text(x + PAD, y + PAD + 17, w - 2 * PAD, 24, [[
        R("Executive-function assessment still leans on ", size=24),
        R("questionnaires (BRIEF-A)", size=24, color=WHITE, bold=True),
        R(" and dry lab drills:", size=24),
    ]])
    bullets(x + PAD, y + PAD + 46, w - 2 * PAD, [
        [R("Stressful & repetitive", size=24, color=WHITE, bold=True), R(" — scores drop with motivation, not ability", size=24)],
        [R("Low engagement", size=24, color=WHITE, bold=True), R(", especially for children and young adults", size=24)],
        [R("Subjective ratings", size=24, color=WHITE, bold=True), R(" — no objective behavioral data", size=24)],
    ], size=24, gap=7)
    rect(x + PAD, y + h - 42, w - 2 * PAD, 30, RGBColor(0x33, 0x24, 0x1A), line=ORANGE, line_w=1.5, radius=0.12)
    text(x + PAD + 6, y + h - 38, w - 2 * PAD - 12, 22, [[
        R("The need: measurement that is objective, repeatable — and actually fun to take.",
          size=22, color=RGBColor(0xFF, 0xD9, 0xB3), bold=True)]], line_spacing=1.15)


def solution(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "02", "The Solution")
    text(x + PAD, y + PAD + 17, w - 2 * PAD, 26, [[
        R("Mission Focus", size=22, color=WHITE, bold=True),
        R(" embeds the assessment inside a 10-minute space mission — each problem above gets a built-in answer:", size=22),
    ]], line_spacing=1.2)
    rows = [
        ("STRESS", 28, [
            R("Assessment is embedded inside ", size=21, color=INK),
            R("meaningful repair tasks", size=21, color=WHITE, bold=True),
            R(" — participants play rather than feel tested, so performance is ", size=21, color=INK),
            R("less distorted by test anxiety", size=21, color=WHITE, bold=True), R(".", size=21, color=INK)]),
        ("BOREDOM", 28, [
            R("A live mission with a real goal — keep the station alive — ", size=21, color=INK),
            R("sustains motivation", size=21, color=WHITE, bold=True),
            R(" for the full session, for kids and adults.", size=21, color=INK)]),
        ("SUBJECTIVE\nDATA", 47, [
            R("The scientific core: a silent ", size=21, color=INK),
            R("behavioral log", size=21, color=WHITE, bold=True),
            R(" of what participants actually do — ", size=21, color=INK),
            R("RT, accuracy, omissions, impulsive actions, task completion", size=21, color=WHITE, bold=True),
            R(" — objective evidence that ", size=21, color=INK),
            R("complements", size=21, color=WHITE, bold=True),
            R(" BRIEF-style ratings and clinical judgment, not replaces them.", size=21, color=INK)]),
        ("PAPERWORK", 19, [
            R("One click → ", size=21, color=INK),
            R("assessor-ready HTML report", size=21, color=WHITE, bold=True),
            R(" + analysis-ready ", size=21, color=INK),
            R("CSV", size=21, color=WHITE, bold=True),
            R(" — no manual scoring.", size=21, color=INK)]),
    ]
    yy = y + PAD + 44
    for tag, rh, runs in rows:
        tag_lines = tag.split("\n")
        tag_h = 6 + 6.5 * len(tag_lines)
        rect(x + PAD, yy + 1, 44, tag_h, RGBColor(0x33, 0x24, 0x1A), line=ORANGE, line_w=1.2, radius=0.18)
        text(x + PAD, yy + 1 + (tag_h - 5.8 * len(tag_lines)) / 2 - 0.6, 44, tag_h,
             [[R(t, font=ORB, size=12, color=RGBColor(0xFF, 0xC0, 0x8A), bold=True)] for t in tag_lines],
             align=PP_ALIGN.CENTER, line_spacing=1.05, space_after=0)
        text(x + PAD + 50, yy, w - 2 * PAD - 50, rh, [runs], line_spacing=1.18)
        yy += rh + 4


def requirements(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "03", "Main Requirements")
    bullets(x + PAD, y + PAD + 18, w - 2 * PAD, [
        [R("Playable 10-minute mission", size=24, color=WHITE, bold=True), R(" — a game, not a quiz", size=24)],
        [R("Measure ", size=24), R("attention, working memory, inhibition & planning", size=24, color=WHITE, bold=True)],
        [R("Every action logged automatically", size=24, color=WHITE, bold=True), R(" during play", size=24)],
        [R("One-click HTML report + CSV", size=24, color=WHITE, bold=True), R(" for the assessor", size=24)],
        [R("All data stays on the device", size=24, color=WHITE, bold=True), R(" — privacy by design", size=24)],
    ], size=24, gap=7, marker="✓", mcolor=GREEN)


def snapshots(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "", "Mission Snapshots", title_size=24)
    cw = (w - 2 * PAD - 8) / 2
    ch = (h - PAD - 18 - PAD - 3 * 14 - 2 * 6) / 3  # 3 rows: cell + caption(14), 6mm row gaps
    from PIL import Image as PILImage
    cells = [
        ("Radar Scan console — attention trials", A("shot_radar.png"), "shot"),
        ("Aboard the station — memorize the code", A("shot_station.png"), "shot"),
        ("Stroop console — match the ink color", A("shot_stroop.png"), "shot"),
        ("Life-Support wires — planning task", A("shot_wires.png"), "shot"),
        ("Your astronaut — playable character", A("char_standing.png"), "char"),
        ("Assessor HTML report — EF profile", A("shot_report.png"), "shot"),
    ]
    for i, (cap, img, kind) in enumerate(cells):
        cx_ = x + PAD + (i % 2) * (cw + 8)
        cy_ = y + PAD + 18 + (i // 2) * (ch + 14 + 6)
        if kind == "shot":
            pic = slide.shapes.add_picture(img, Mm(cx_), Mm(cy_), Mm(cw), Mm(ch))
            iw, ih = PILImage.open(img).size  # crop source to cell aspect (center)
            src_ar, cell_ar = iw / ih, cw / ch
            if src_ar > cell_ar:
                cut = (1 - cell_ar / src_ar) / 2
                pic.crop_left = cut; pic.crop_right = cut
            else:
                cut = (1 - src_ar / cell_ar) / 2
                pic.crop_top = cut; pic.crop_bottom = cut
            pic.line.color.rgb = EDGE; pic.line.width = Pt(1)
        elif kind == "char":
            rect(cx_, cy_, cw, ch, RGBColor(0x10, 0x22, 0x44), line=EDGE, radius=0.06)
            iw, ih = PILImage.open(img).size
            img_h = ch - 4
            img_w = img_h * iw / ih
            slide.shapes.add_picture(img, Mm(cx_ + (cw - img_w) / 2), Mm(cy_ + ch - img_h - 1),
                                     Mm(img_w), Mm(img_h))
        else:
            rect(cx_, cy_, cw, ch, RGBColor(0x09, 0x14, 0x2C), line=CYANSOFT, line_w=1.4,
                 radius=0.06, dash=MSO_LINE_DASH_STYLE.DASH)
            text(cx_, cy_ + ch / 2 - 5, cw, 10, [[R("SCREENSHOT SLOT", font=ORB, size=12, color=CYANSOFT)]],
                 align=PP_ALIGN.CENTER)
        text(cx_, cy_ + ch + 2, cw, 12, [[R(cap, size=17, color=SOFT)]], line_spacing=1.1)


# ---------- CENTER ----------
def sys_flow(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "04", "System Flow")
    steps = [
        ("1", "Participant Onboarding", [R("ID and age entered on a station terminal", size=23)]),
        ("2", "Interactive Tutorial", [R("Movement & console training", size=23)]),
        ("3", "The Mission — 10 minutes", [R("Repair tasks appear across ", size=23), R("4 stations", size=23, color=WHITE, bold=True),
                                           R(" — the player chooses what to handle first", size=23)]),
        ("4", "Silent Telemetry", [R("Every reaction recorded in the background", size=23)]),
        ("5", "Assessor Report", [R("Executive-function profile + CSV", size=23, color=WHITE, bold=True),
                                  R(", generated instantly", size=23)]),
    ]
    sy = y + PAD + 20
    row_h = (h - PAD - 24 - PAD) / 5
    for i, (n, t, drs) in enumerate(steps):
        yy = sy + i * row_h
        if i < 4:
            ln = rect(x + PAD + 9, yy + 20, 1, row_h - 18, CYAN)
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Mm(x + PAD), Mm(yy), Mm(19), Mm(19))
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x14, 0x30, 0x5C)
        c.line.color.rgb = CYAN; c.line.width = Pt(2); c.shadow.inherit = False
        text(x + PAD, yy + 3.2, 19, 12, [[R(n, font=ORB, size=21, color=CYAN, bold=True)]], align=PP_ALIGN.CENTER)
        text(x + PAD + 26, yy - 1, w - 2 * PAD - 26, 12, [[R(t, size=25, color=WHITE, bold=True)]])
        text(x + PAD + 26, yy + 10.5, w - 2 * PAD - 26, row_h - 10, [drs], line_spacing=1.15)


def tasks(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "05", "Cognitive Tasks")
    data = [
        ("RADAR SCAN", "NAVIGATION STATION", "ATTENTION · TASK-MONITOR", CYAN,
         [R("Spot ", size=22), R("rare asteroid contacts", size=22, color=WHITE, bold=True),
          R(" among streams of debris.", size=22)],
         "Hits · False alarms · d′ sensitivity"),
        ("CODE RECALL", "ENGINE STATION", "WORKING MEMORY", VIOLET,
         [R("Memorize a ", size=22), R("4-digit reactor code", size=22, color=WHITE, bold=True),
          R(", recall it across the station.", size=22)],
         "Recall accuracy · Typos · Timeout"),
        ("STROOP CONSOLE", "COMMS STATION", "INHIBITION", ORANGE,
         [R("The word and its ink color disagree — answer the ", size=22), R("color", size=22, color=WHITE, bold=True),
          R(", not the word.", size=22)],
         "RT mean/SD · Commissions · Post-error slowing"),
        ("POWER CELL RUN", "LIFE-SUPPORT STATION", "PLAN · ORGANIZE", GREEN,
         [R("Find a power cell, route it across the station, ", size=22),
          R("re-wire the socket", size=22, color=WHITE, bold=True), R(" in time.", size=22)],
         "Pickup latency · Wiring errors · Completion time"),
    ]
    cw = (w - 2 * PAD - 9) / 2
    ch = (h - PAD - 22 - 9 - PAD) / 2
    for i, (t, st, dom, dc, desc, chips) in enumerate(data):
        cx_ = x + PAD + (i % 2) * (cw + 9)
        cy_ = y + PAD + 18 + (i // 2) * (ch + 9)
        rect(cx_, cy_, cw, ch, CARD2, line=EDGE, radius=0.055)
        text(cx_ + 8, cy_ + 7, cw - 16, 12, [[R(t, font=ORB, size=20, color=WHITE, bold=True)]])
        text(cx_ + 8, cy_ + 17.5, cw - 16, 9, [[R(st, size=16, color=SOFT, bold=True)]])
        pill_w = 6 + len(dom) * 3.5
        rect(cx_ + 8, cy_ + 26.5, pill_w, 11, None, line=dc, line_w=1.2, radius=0.5)
        text(cx_ + 8, cy_ + 28.6, pill_w, 7, [[R(dom, size=15, color=dc, bold=True)]], align=PP_ALIGN.CENTER)
        text(cx_ + 8, cy_ + 42, cw - 16, ch - 60, [desc], line_spacing=1.22)
        text(cx_ + 8, cy_ + ch - 13, cw - 16, 10, [[R(chips, size=16, color=RGBColor(0xEE, 0xF6, 0xFF), bold=True)]])


def pipeline(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "", "Under the Hood — Measurement Pipeline", title_size=22)
    rows = [
        ("Mission console", "runs timed trials"),
        ("Telemetry", "every action time-stamped"),
        ("Metrics", "accuracy, speed, errors per task"),
        ("Report", "HTML profile · CSV"),
    ]
    rh = (h - PAD - 20 - PAD - 3 * 5) / 4
    yy = y + PAD + 18
    for name, desc in rows:
        rect(x + PAD + 6, yy, w - 2 * PAD - 6, rh, RGBColor(0x0A, 0x15, 0x2E), line=EDGE, radius=0.18)
        text(x + PAD + 13, yy + rh / 2 - 5.5, w - 2 * PAD - 20, 10, [[
            R(name, font=MONO, size=21, color=WHITE, bold=True),
            R("   " + desc, size=19, color=SOFT),
        ]])
        yy += rh + 5
    rect(x + PAD, y + PAD + 20, 1.2, h - 2 * PAD - 22, ORANGE)


# ---------- RIGHT ----------
def tech(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "06", "Technologies & Tools")
    rows = [
        ("Unity 6 · URP", "3D engine & rendering pipeline", CYAN),
        ("C#", "gameplay & assessment logic", VIOLET),
        ("Unity Input System", "first-person controls", GREEN),
        ("TextMeshPro", "HUD & console UI", ORANGE),
        ("ElevenLabs", "AI mission voice-overs", CYAN),
        ("Meshy AI", "3D station props & consoles", VIOLET),
        ("Git · GitHub", "version control & delivery", GREEN),
    ]
    rh = (h - PAD - 20 - PAD - 6 * 3.5) / 7
    yy = y + PAD + 18
    for name, desc, dc in rows:
        rect(x + PAD, yy, w - 2 * PAD, rh, RGBColor(0x0A, 0x15, 0x2E), line=EDGE, radius=0.22)
        sq = rect(x + PAD + 5, yy + rh / 2 - 3.4, 6.8, 6.8, dc, radius=0.3)
        text(x + PAD + 16, yy + rh / 2 - 5.5, w - 2 * PAD - 22, 10, [[
            R(name, size=21, color=WHITE, bold=True),
            R("  — " + desc, size=19, color=SOFT),
        ]])
        yy += rh + 3.5


def testing(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "07", "Testing & Evaluation")
    bullets(x + PAD, y + PAD + 18, w - 2 * PAD, [
        [R("Rapid test missions", size=22, color=WHITE, bold=True), R(" — full pipeline checked end-to-end", size=22)],
        [R("Live assessor view", size=22, color=WHITE, bold=True), R(" verified during real missions", size=22)],
        [R("Pilot playtests", size=22, color=WHITE, bold=True), R(" with peers ahead of the exhibition", size=22)],
    ], size=22, gap=7)


def challenges(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "08", "Challenges & Insights")
    items = [
        ("01", [R("Fun vs. validity", size=22, color=WHITE, bold=True), R(" — free play outside, strict timing inside every console", size=22)]),
        ("02", [R("Invisible measurement", size=22, color=WHITE, bold=True), R(" — the “test” stays out of sight", size=22)]),
        ("03", [R("One shared core", size=22, color=WHITE, bold=True), R(" — a new cognitive task is a single script", size=22)]),
        ("04", [R("AI assets", size=22, color=WHITE, bold=True), R(" — generated 3D props needed Unity rework", size=22)]),
    ]
    paras = [[R(n + "  ", font=ORB, size=18, color=ORANGE, bold=True)] + rs for n, rs in items]
    text(x + PAD, y + PAD + 18, w - 2 * PAD, h - 70, paras, space_after=7, line_spacing=1.22)
    rect(x + PAD, y + h - 42, w - 2 * PAD, 30, RGBColor(0x10, 0x2A, 0x44), line=CYAN, line_w=1.5, radius=0.1)
    text(x + PAD + 6, y + h - 38, w - 2 * PAD - 12, 22, [[
        R("Players forget they are being measured — and strict timing is exactly what keeps the data meaningful.",
          size=20, color=INK, italic=True)]], line_spacing=1.15)


def review(x, y, w, h):
    sec_title(x + PAD, y + PAD, w, "09", "Project Review")
    items = [
        ("✔", GREEN, [R("Goals met", size=22, color=WHITE, bold=True), R(" — all four cognitive domains playable end-to-end", size=22)]),
        ("✔", GREEN, [R("~26 metrics per session", size=22, color=WHITE, bold=True), R(" — HTML + CSV + raw log", size=22)]),
        ("✔", GREEN, [R("Demo-ready", size=22, color=WHITE, bold=True), R(" — stable 10-minute missions", size=22)]),
        ("↻", ORANGE, [R("In hindsight", size=22, color=WHITE, bold=True), R(" — playtest earlier", size=22)]),
        ("→", CYAN, [R("Next", size=22, color=WHITE, bold=True), R(" — validation study against standard BRIEF-A scores", size=22)]),
    ]
    paras = [[R(m + "  ", size=22, color=c, bold=True)] + rs for m, c, rs in items]
    text(x + PAD, y + PAD + 18, w - 2 * PAD, h - 40, paras, space_after=7, line_spacing=1.22)


def takeaway(x, y, w, h):
    # gradient-ish border via outer cyan rect
    rect(x - 1, y - 1, w + 2, h + 2, CYAN, radius=0.06)
    rect(x + 0.6, y + 0.6, w - 1.2, h - 1.2, RGBColor(0x0B, 0x18, 0x34), radius=0.06)
    sec_title(x + PAD, y + PAD, w, "10", "Key Takeaway")
    text(x + PAD, y + PAD + 18, w - 2 * PAD, h - 40, [[
        R("Assessment doesn't have to feel like a test. A ", size=24.5, color=INK),
        R("10-minute space mission", size=24.5, color=ORANGE, bold=True),
        R(" can quietly produce a ", size=24.5, color=INK),
        R("full executive-function profile", size=24.5, color=CYAN, bold=True),
        R(" — engagement and measurement in the same loop.", size=24.5, color=INK),
    ]], line_spacing=1.3)


def qr_row(x, y, w, h):
    rect(x + PAD - 4, y + 6, 48, 48, WHITE, radius=0.1)
    slide.shapes.add_picture(A("qr_github.png"), Mm(x + PAD - 1), Mm(y + 9), Mm(42), Mm(42))
    text(x + PAD + 50, y + 10, w - PAD - 60, 12, [[R("FULL CODE & DOCS", font=ORB, size=18, color=WHITE, bold=True)]])
    text(x + PAD + 50, y + 22, w - PAD - 60, 12, [[R("github.com/lidorbenhamo26/Final-Project", size=17, color=WHITE, bold=True)]])
    text(x + PAD + 50, y + 32, w - PAD - 60, 12, [[R("Unity project · project book · demo video", size=16, color=SOFT)]])


# ---------- build columns ----------
stack(LX, LW, [
    (165, bg_need),
    (195, solution),
    (130, requirements),
    (328, snapshots),
])
stack(CX, CW, [
    (305, sys_flow),
    (330, tasks),
    (168, pipeline),
])
stack(RX, RW, [
    (185, tech),
    (115, testing),
    (185, challenges),
    (140, review),
    (110, takeaway),
    (60, qr_row),
])

# ============ STATS BAND ============
BY = 1056
rect(22, BY, 797, 68, CARD, line=EDGE, radius=0.08)
stats = [("4", "COGNITIVE STATIONS", CYAN), ("10 MIN", "ONE FULL MISSION", ORANGE),
         ("~26", "METRICS PER SESSION", VIOLET), ("3", "EXPORT ARTIFACTS", GREEN)]
cw_ = 150
for i, (v, l, c) in enumerate(stats):
    sx = 40 + i * cw_
    text(sx, BY + 10, cw_, 24, [[R(v, font=ORB, size=44, color=c, bold=True)]], align=PP_ALIGN.CENTER)
    text(sx, BY + 40, cw_, 12, [[R(l, size=18, color=SOFT, bold=True)]], align=PP_ALIGN.CENTER)
    if i:
        rect(sx - 2, BY + 12, 0.5, 44, EDGE)
# wave astronaut + bubble
slide.shapes.add_picture(A("char_wave.png"), Mm(742), Mm(BY - 14), height=Mm(80))
rect(648, BY + 6, 88, 15, RGBColor(0x10, 0x2A, 0x44), line=CYANSOFT, radius=0.5)
text(648, BY + 9.4, 88, 9, [[R("Ready for launch!", size=16, color=CYANSOFT, bold=True)]], align=PP_ALIGN.CENTER)

# ============ FOOTER ============
rect(0, 1168, 841, 21, RGBColor(0x05, 0x0B, 0x1A), line=None)
rect(0, 1168, 841, 0.8, EDGE)
text(0, 1173.5, 841, 12, [[
    R("Braude College of Engineering, Karmiel", size=17.5, color=WHITE, bold=True),
    R("   ◆   Department of Software Engineering & Information Systems   ◆   Capstone Project 61999 — Phase B   ◆   ", size=17.5, color=SOFT),
    R("2026", size=17.5, color=WHITE, bold=True),
]], align=PP_ALIGN.CENTER)

out = os.path.join(HERE, "MissionFocus_Poster_A0.pptx")
prs.save(out)
print("saved", out)
